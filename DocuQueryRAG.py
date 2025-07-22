import gradio as gr
import numpy as np
import faiss
import argparse
from sentence_transformers import SentenceTransformer
import subprocess
import pdfplumber
import re
from docx import Document
from pptx import Presentation

# --- Parse CLI arguments ---
parser = argparse.ArgumentParser(description="RAG DocQA with custom LLM support via Ollama")
parser.add_argument('--model', type=str, default='mistral:7b', help='Model name for Ollama (e.g., mistral:7b, llama3:8b)')
args = parser.parse_args()
selected_model = args.model

# --- Utilities for Chunking ---
def simple_sent_tokenize(text):
    return re.split(r'(?<=[.!?]) +', text)

def chunk_text(text, max_words=200):
    sentences = simple_sent_tokenize(text)
    chunks, current_chunk, current_words = [], [], 0
    for sentence in sentences:
        word_count = len(sentence.split())
        if current_words + word_count > max_words:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentence]
            current_words = word_count
        else:
            current_chunk.append(sentence)
            current_words += word_count
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    return chunks

# --- File Extractors ---
def extract_text_from_pdf(file_obj):
    text = ""
    with pdfplumber.open(file_obj) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_text_from_docx(file_obj):
    return "\n".join([para.text for para in Document(file_obj.name).paragraphs if para.text.strip()])

def extract_text_from_pptx(file_obj):
    prs = Presentation(file_obj.name)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_txt(file_obj):
    return file_obj.read().decode("utf-8")

def extract_text_from_file(file_obj):
    name = file_obj.name.lower()
    try:
        if name.endswith(".pdf"):
            return extract_text_from_pdf(file_obj)
        elif name.endswith(".docx"):
            return extract_text_from_docx(file_obj)
        elif name.endswith(".pptx"):
            return extract_text_from_pptx(file_obj)
        elif name.endswith(".txt"):
            return extract_text_from_txt(file_obj)
        else:
            return ""
    except Exception as e:
        return f"[ERROR: failed to read {name}] {str(e)}"

# --- Embedding + FAISS ---
embedder = SentenceTransformer('all-MiniLM-L6-v2')

def embed_chunks(chunks):
    return np.array(embedder.encode(chunks, show_progress_bar=False))

def build_faiss_index(embeddings):
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(embeddings)
    return index

def search_index(query, embedder, index, chunks, k=5):
    query_vector = embedder.encode([query])
    distances, indices = index.search(np.array(query_vector), k)
    return [chunks[i] for i in indices[0]]

# --- LLM via Ollama ---
def call_ollama_llm(prompt, model='mistral:7b'):
    try:
        result = subprocess.run(
            ['ollama', 'run', model],
            input=prompt.encode('utf-8'),
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
        return result.stdout.decode('utf-8').strip()
    except Exception as e:
        return f"[ERROR] Failed to call Ollama: {e}"

def generate_answer(query, top_chunks, model='mistral:7b'):
    context = "\n\n".join(top_chunks)
    prompt = f"""
You are an expert assistant. Use the following context to answer the question.

Context:
{context}

Question:
{query}

Answer:
""".strip()
    return call_ollama_llm(prompt, model)

# --- Global State ---
stored_chunks = []
stored_index = None
chat_history = []

# --- Upload Handler ---
def upload_files(files):
    global stored_chunks, stored_index, chat_history
    all_text = ""
    for file_obj in files:
        text = extract_text_from_file(file_obj)
        if text and not text.startswith("[ERROR"):
            all_text += text + "\n"

    if not all_text.strip():
        return "❌ No valid text found."

    chunks = chunk_text(all_text)
    embeddings = embed_chunks(chunks)
    index = build_faiss_index(embeddings)
    stored_chunks = chunks
    stored_index = index
    chat_history.clear()
    return f"✅ {len(files)} files processed. Total Chunks: {len(chunks)}"

# --- Query Handler ---
def query_doc(question, history):
    if not stored_chunks or not stored_index:
        return history + [["User", question], ["System", "❌ Please upload documents first."]]

    top_chunks = search_index(question, embedder, stored_index, stored_chunks)
    answer = generate_answer(question, top_chunks, model=selected_model)
    formatted_answer = "\n".join(answer.split(". "))  # Break long lines
    history.append(["User", question])
    history.append(["System", formatted_answer])
    return history

# --- Reset Memory ---
def reset_all():
    global stored_chunks, stored_index, chat_history
    stored_chunks = []
    stored_index = None
    chat_history = []
    return [], "🧹 Cleared memory. You can now upload new documents."

# --- Gradio App ---
with gr.Blocks(theme=gr.themes.Base()) as demo:
    gr.Markdown("## 📚 RAG Chatbot with Multi-File Upload and FAISS")

    with gr.Row():
        file_uploader = gr.File(file_types=[".pdf", ".docx", ".pptx", ".txt"], file_count="multiple", label="📂 Upload Documents")
        upload_btn = gr.Button("📥 Process Files")
    upload_output = gr.Textbox(label="Status", interactive=False)

    chatbot = gr.Chatbot(label="💬 Chat History", height=200)
    question_box = gr.Textbox(lines=1, placeholder="Ask a question and press Enter or click the button...")
    ask_btn = gr.Button("🔍 Get Answer")
    reset_btn = gr.Button("🧹 Reset")

    upload_btn.click(upload_files, inputs=[file_uploader], outputs=[upload_output])
    ask_btn.click(query_doc, inputs=[question_box, chatbot], outputs=[chatbot])
    question_box.submit(query_doc, inputs=[question_box, chatbot], outputs=[chatbot])
    reset_btn.click(reset_all, outputs=[chatbot, upload_output])

demo.launch()
